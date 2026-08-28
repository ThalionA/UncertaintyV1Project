# -*- coding: utf-8 -*-
"""Build the lean, figure-first meeting deck (one figure per slide + a short
title and a one-line takeaway). One-off; lives in the gitignored figures/ dir."""
from pathlib import Path
from PIL import Image
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

FIGDIR = Path("/Users/theoamvr/Documents/ResearchVault/attachments/2026-06-18-meeting")
OUT = Path("/Users/theoamvr/Desktop/Experiments/UncertaintyV1/nn_decoder/figures/meeting_2026_06_18_deck.pptx")

NAVY = RGBColor(0x1F, 0x2D, 0x3D)
GREY = RGBColor(0x55, 0x55, 0x55)
ACCENT = RGBColor(0xE6, 0x55, 0x0D)   # PCA_EVAR orange — the through-line colour
WHITE = RGBColor(0xFF, 0xFF, 0xFF)

SLIDES = [
    ("Sharpen / broaden & location — on the real IO posteriors",
     "On the real perceptual posteriors: every loss penalises a location shift equally (the shared signal). "
     "KL/JS/CE penalise sharpening more than broadening (a restoring force); Projection-based & Wasserstein are symmetric.",
     "mtg0618_locsharp_sweeps.png"),
    ("Real trained decoders vs the IO target — on real trials",
     "Overlay each loss-trained decoder's real decoded posterior on the IO target: the projection-based & Wasserstein "
     "decoders over-sharpen (jagged spikes), even on broad/bimodal targets; CE/KL/JS track the shape. No fitting.",
     "mtg0618_locsharp_examples.png"),
    ("Projection-based's mismatch along each PC — real V1 (6 mice)",
     "All losses match the location subspace (PC0-1); Projection-based's shape-subspace error is 38x KL's. "
     "The over-sharpening lives in the trailing high-frequency 'shape' PCs.",
     "mtg0618_pc_mismatch.png"),
    ("Spatial vs temporal — train objective x eval metric, with & without Mouse 2",
     "Normalised-loss difference (spat-temp). Projection-based: a wash under its OWN metric (Δ+0.05); spatial >> temporal "
     "under KL (Δ-0.83 all 6 mice; -0.91 dropping M2) and JS. Robust to M2; only a calibrated metric sees it.",
     "mtg0618_spat_temp_crossloss.png"),
    ("Spatial vs temporal — per animal (with & without Mouse 2) + neuron count",
     "Projection-based is the only loss where spatial beats temporal (normalised KL loss 1.34 vs 2.17, p=0.01); robust to "
     "dropping Mouse 2 (Δ-0.91); neuron count (65-153) does not predict performance (n=6).",
     "mtg0618_spat_temp_per_animal.png"),
    ("Dropout vs early stopping",
     "Early-stop halves Projection-based peakiness (0.72->0.36) and raw KL (4.6->1.75); dropout does nothing (even slightly "
     "worse). The Projection-based metric is blind to all of it — the remedy is a width term in the loss.",
     "mtg0618_dropout_vs_earlystop.png"),
    ("Train-val gap with dropout (the monitor_val view)",
     "A static offset, not progressive overfitting — val loss plateaus by epoch ~20 and never climbs; "
     "dropout barely closes it (spatial ~12%, temporal flat). Gap != over-sharpening.",
     "mtg0618_dropout_trainval.png"),
    ("Do the temporal bins sample? Location & width similarity",
     "Bins are broad copies: mean per-bin width ≈ target (19°), bin-to-bin location spread ~15° (below target), "
     "~flat across λ_H. Only Projection-based sharpens its bins (20°->15°) — but locations don't spread to compensate.",
     "mtg0618_bin_similarity.png"),
    ("Per-bin examples — twin axis (target left; average & bins right)",
     "IO target on the LEFT axis; the time-average AND the 10 per-bin posteriors on the RIGHT axis at full height. "
     "Projection-based/Wasserstein bins are spiky; CE/KL/JS bins are broad ≈ the target.",
     "mtg0618_bin_examples.png"),
    ("Does the bins' similarity depend on the stimulus?",
     "Between-bin dispersion grouped by stimulus: location dispersion RISES with stimulus dispersion (bins spread more "
     "when uncertain); width dispersion is U-shaped in orientation (largest at the 0/90° references).",
     "mtg0618_bin_by_condition.png"),
    ("Peakiness vs uncertainty",
     "Over-confidence grows where the ideal observer is least certain — Projection-based up to ~5x, peaking at the 45° "
     "boundary. Caveat: structured over-sharpening, not 'ignoring the boundary' (raw peakiness dips there too).",
     "mtg0618_peakiness_vs_uncertainty.png"),
    ("Shuffle control — three 'no trial-information' nulls",
     "Under KL, trained Projection-based 2.2x/3.4x and Wasserstein 2.7x/2.0x WORSE than predicting the mean; CE/KL/JS beat "
     "all three nulls. Projection-based/Wasserstein 'win' only on their own metric. Lead with predict-mean (strictest).",
     "mtg0618_shuffle_nulls.png"),
]


def _txt(slide, l, t, w, h, text, size, color, bold=False, italic=False,
         align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run(); r.text = text
    f = r.font
    f.size = Pt(size); f.bold = bold; f.italic = italic; f.color.rgb = color
    f.name = "Calibri"
    return tb


def _bar(slide, color, h=0.13):
    from pptx.enum.shapes import MSO_SHAPE
    sh = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0),
                                Inches(13.333), Inches(h))
    sh.fill.solid(); sh.fill.fore_color.rgb = color
    sh.line.fill.background()
    sh.shadow.inherit = False
    return sh


def _fit(slide, path, box_l, box_t, box_w, box_h):
    iw, ih = Image.open(path).size
    scale = min(box_w / iw, box_h / ih)
    w, h = iw * scale, ih * scale
    left = box_l + (box_w - w) / 2
    top = box_t + (box_h - h) / 2
    slide.shapes.add_picture(str(path), Inches(left), Inches(top),
                             Inches(w), Inches(h))


prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# --- title slide ---
s = prs.slides.add_slide(BLANK)
_bar(s, ACCENT, h=0.22)
_txt(s, 0.8, 2.25, 11.7, 2.2,
     "V1 uncertainty: loss geometry, spatial-temporal,\ndropout & temporal sampling",
     38, NAVY, bold=True)
_txt(s, 0.8, 4.5, 11.7, 1.6,
     "Figures for the 2026-06-18 meeting  ·  loss_comparison_v1 (6 mice, Q / half / 100 ms) + "
     "lambdaH_sweep\nLosses: Projection-based · CE · KL · JS · Wasserstein   |   companion report in ResearchVault",
     17, GREY)

# --- content slides ---
for title, takeaway, fig in SLIDES:
    s = prs.slides.add_slide(BLANK)
    _bar(s, ACCENT)
    _txt(s, 0.4, 0.22, 12.53, 0.9, title, 22, NAVY, bold=True)
    # the takeaway goes in the SPEAKER NOTES, never as a box over the figure
    s.notes_slide.notes_text_frame.text = takeaway
    p = FIGDIR / fig
    if p.exists():
        _fit(s, p, 0.4, 1.2, 12.53, 6.1)
    else:
        _txt(s, 0.4, 3.0, 12.53, 1.0, f"[missing figure: {fig}]", 14, ACCENT)

prs.save(str(OUT))
print(f"Saved {OUT}  ({len(prs.slides)} slides)")
